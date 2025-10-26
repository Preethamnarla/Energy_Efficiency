# build_simple_ppt_annotated.py
# Generates a clean, visuals-first deck:
#   - Energy_Efficient_Simulator_Simple.pptx
#   - assets/ images for visuals (including annotated topology)

import os, math
import numpy as np
import matplotlib.pyplot as plt
import networkx as nx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = "Energy_Efficient_Simulator_Simple.pptx"
ASSETS = "assets"
os.makedirs(ASSETS, exist_ok=True)

# ---------------- Theme (light) ----------------
COL_BG = RGBColor(255, 255, 255)
COL_TITLE = RGBColor(15, 23, 42)
COL_TEXT = RGBColor(51, 65, 85)
FONT = "Arial"

# -------------- Helpers for images -------------
def fig_save(fig, path):
    fig.tight_layout()
    fig.savefig(path, dpi=200, transparent=True)
    plt.close(fig)

def pos_hex(R=1.0):
    pos = {0: (0.0, 0.0)}
    ang = np.linspace(0, 2*np.pi, 6, endpoint=False)
    for i,a in enumerate(ang, start=1):
        pos[i] = (R*math.cos(a), R*math.sin(a))
    return pos

# -------------- Visuals: Annotated Topology --------------
def draw_topology_annotated(path):
    import matplotlib.patches as patches
    plt.rcParams.update({"font.size": 10})
    fig, ax = plt.subplots(figsize=(7.4, 6.6))

    G = nx.Graph(); pos = pos_hex(1.0)
    for n in pos: G.add_node(n)
    for i in range(1, 7): G.add_edge(0, i)

    # Edges (neighbor/interference links)
    nx.draw_networkx_edges(G, pos, ax=ax, width=2, edge_color="#cbd5e1", alpha=0.85)

    # Nodes: center vs ring
    node_colors = {0: "#f59e0b"}  # center = amber
    for i in range(1, 7): node_colors[i] = "#3b82f6"  # ring = blue
    nx.draw_networkx_nodes(
        G, pos, ax=ax, node_size=1300,
        node_color=[node_colors[n] for n in G.nodes()],
        edgecolors="#0f172a", linewidths=1.6
    )
    nx.draw_networkx_labels(
        G, pos, ax=ax, labels={n: f"Site {n}" for n in pos},
        font_size=10, font_weight="bold", font_color="white"
    )

    # Mobility arrows (outer ring circulation)
    ring = [1, 2, 3, 4, 5, 6, 1]
    for a, b in zip(ring[:-1], ring[1:]):
        (x1, y1), (x2, y2) = pos[a], pos[b]
        ax.annotate(
            "", xy=(x2, y2), xytext=(x1, y1),
            arrowprops=dict(arrowstyle="-|>", color="#64748b", lw=1.6, alpha=0.6,
                            connectionstyle="arc3,rad=0.25")
        )

    # Callouts
    x0, y0 = pos[0]
    ax.text(x0, y0 - 0.38, "Hub site (higher activity)", ha="center", va="top",
            fontsize=9, color="#0f172a",
            bbox=dict(boxstyle="round,pad=0.3", facecolor="white", edgecolor="#cbd5e1", alpha=0.97))

    x2, y2 = pos[2]
    ax.text(x2 + 0.72, y2 + 0.32, "Outer ring: more sleep opportunities\nwhen demand is low",
            ha="left", va="center", fontsize=9, color="#0f172a",
            bbox=dict(boxstyle="round,pad=0.3", facecolor="white", edgecolor="#cbd5e1", alpha=0.97))

    # Legend (concise config)
    legend_lines = [
        "Scenario: 7‑site hex (UMa), reuse‑1",
        "Each site: 3 sectors (65° panels)",
        "Carrier/BW: ~3.5 GHz / 20–40 MHz",
        "Traffic: diurnal; mobility across cells",
        "Interference: neighbor coupling on",
        "Control: baselines vs RL (sleep/power)"
    ]
    ax.text(1.02, 0.98, "Topology & Assumptions\n" + "\n".join("• " + s for s in legend_lines),
            transform=ax.transAxes, va="top", ha="left", fontsize=9, color="#111827",
            bbox=dict(boxstyle="round,pad=0.45", facecolor="white", edgecolor="#cbd5e1", alpha=0.98))

    # Sectorization inset (3 sectors/site)
    inset = ax.inset_axes([0.79, 0.06, 0.18, 0.18])  # x, y, w, h in axes fraction
    inset.axis("off")
    center = (0.5, 0.5); R = 0.42
    for k, c in [(0, "#60a5fa"), (120, "#34d399"), (240, "#fbbf24")]:
        wedge = patches.Wedge(center, R, k, k + 80, facecolor=c, edgecolor="#1f2937", lw=1)
        inset.add_patch(wedge)
    inset.text(0.5, 0.06, "3 sectors/site", ha="center", va="bottom", fontsize=8, color="#334155")

    ax.set_title("7‑Site Hex Topology (mobility, interference, sectorization)", fontsize=12, pad=10)
    ax.set_aspect("equal"); ax.axis("off")
    fig_save(fig, path)

# -------------- Visuals: Architecture --------------
def draw_architecture(path):
    import matplotlib.patches as patches
    fig, ax = plt.subplots(figsize=(10.2, 4.6))
    ax.axis("off")

    def box(x,y,w,h,text,color):
        r = patches.FancyBboxPatch((x,y), w,h, boxstyle="round,pad=0.02,rounding_size=0.03",
                                   linewidth=2, edgecolor=color, facecolor="white")
        ax.add_patch(r)
        ax.text(x+w/2, y+h/2, text, ha="center", va="center", fontsize=10, color="#1f2937")

    # Top flow
    box(0.03,0.56,0.18,0.25,"Topology\n(NetworkX)", "#60a5fa")
    box(0.23,0.56,0.18,0.25,"Traffic & Mobility", "#60a5fa")
    box(0.43,0.56,0.18,0.25,"Channel Models\nPath loss • Shadowing • Fading", "#60a5fa")
    box(0.63,0.56,0.18,0.25,"PHY/MAC\n(CQI → MCS, Schedulers)", "#60a5fa")
    box(0.83,0.56,0.14,0.25,"Metrics\n& Logs", "#60a5fa")

    # Bottom line
    box(0.28,0.16,0.22,0.25,"Energy Model\nActive • Micro • Deep sleep", "#f472b6")
    box(0.56,0.16,0.22,0.25,"Control Policies\nBaselines • RL Agent", "#34d399")

    # Arrows
    def arrow(x1,y1,x2,y2):
        ax.annotate("", xy=(x2,y2), xytext=(x1,y1),
                    arrowprops=dict(arrowstyle="->", lw=2, color="#94a3b8"))

    arrow(0.21,0.69,0.23,0.69)
    arrow(0.41,0.69,0.43,0.69)
    arrow(0.61,0.69,0.63,0.69)
    arrow(0.81,0.69,0.83,0.69)
    arrow(0.54,0.56,0.67,0.41)  # PHY→Control
    arrow(0.39,0.41,0.56,0.41)  # Energy→Control
    arrow(0.67,0.41,0.73,0.56)  # Control→Metrics

    ax.set_title("Simulator Architecture (modules & flow)", fontsize=12, pad=10)
    fig_save(fig, path)

# -------------- Visuals: RL Loop --------------
def draw_rl_loop(path):
    import matplotlib.patches as patches
    fig, ax = plt.subplots(figsize=(8.2, 3.6))
    ax.axis("off")

    def rounded(x,y,w,h,label,color):
        r = patches.FancyBboxPatch((x,y), w,h, boxstyle="round,pad=0.02,rounding_size=0.03",
                                   linewidth=2, edgecolor=color, facecolor="white")
        ax.add_patch(r)
        ax.text(x+w/2, y+h/2, label, ha="center", va="center", fontsize=11, color="#111827")

    rounded(0.05,0.30,0.22,0.38,"Environment\n(Simulator)", "#60a5fa")
    rounded(0.73,0.30,0.22,0.38,"RL Agent\n(Policy)", "#34d399")

    def arrow(txt, x1,y1,x2,y2, above=True):
        ax.annotate("", xy=(x2,y2), xytext=(x1,y1), arrowprops=dict(arrowstyle="->", lw=2, color="#94a3b8"))
        ax.text((x1+x2)/2, (y1+y2)/2 + (0.07 if above else -0.09), txt, fontsize=10, color="#334155", ha="center")

    arrow("State: load, CQI, queues, power state", 0.27,0.49,0.73,0.49, above=True)
    arrow("Action: sleep/power, allocation hints", 0.73,0.43,0.27,0.43, above=False)
    ax.text(0.5,0.16,"Reward prioritizes energy efficiency with QoS guardrails", ha="center", fontsize=10, color="#334155")

    ax.set_title("Reinforcement learning loop (high‑level)", fontsize=12, pad=10)
    fig_save(fig, path)

# -------------- Visuals: Energy States --------------
def draw_energy_states(path):
    import matplotlib.patches as patches
    fig, ax = plt.subplots(figsize=(8.0, 3.1))
    ax.axis("off")

    def state(x,label,color):
        r = patches.FancyBboxPatch((x,0.35), 0.22,0.30, boxstyle="round,pad=0.02,rounding_size=0.03",
                                   linewidth=2, edgecolor=color, facecolor="white")
        ax.add_patch(r)
        ax.text(x+0.11,0.50,label, ha="center", va="center", fontsize=10, color="#111827")

    state(0.10,"Active", "#3b82f6")
    state(0.39,"Micro‑sleep\n(fast wake)", "#10b981")
    state(0.68,"Deep sleep\n(max save)", "#f59e0b")

    def arrow(x1,x2):
        ax.annotate("", xy=(x2,0.50), xytext=(x1,0.50),
                    arrowprops=dict(arrowstyle="<->", lw=2, color="#94a3b8"))

    arrow(0.32,0.39); arrow(0.61,0.68)
    ax.text(0.50,0.20,"Guard timers prevent frequent toggling; QoS checks before deeper sleep.",
            ha="center", fontsize=10, color="#334155")
    ax.set_title("Energy states & transitions", fontsize=12, pad=10)
    fig_save(fig, path)

# -------------- Visuals: Dashboard Mock --------------
def draw_dashboard_mock(path):
    np.random.seed(0)
    t = np.linspace(0,24,96)
    load = 0.2 + 0.6*(np.sin((t-8)/24*2*np.pi)+1)/2
    base_power = 700 + 500*load
    rl_power = base_power*0.9

    fig, axs = plt.subplots(1,3, figsize=(10.2,3.3))

    # 1) Time series
    axs[0].plot(t, base_power, color="#ef4444", lw=2, label="Baseline")
    axs[0].plot(t, rl_power, color="#10b981", lw=2, label="RL")
    axs[0].set_title("Power vs Time"); axs[0].set_xlabel("Hours"); axs[0].set_yticks([]); axs[0].grid(alpha=0.2)
    axs[0].legend(frameon=False, fontsize=8)

    # 2) Stacked bars (no precise numbers)
    idle = [7,5]; dyn = [5,5]; wake = [0.0,0.4]
    labels = ["Baseline","RL"]
    axs[1].bar(labels, idle, color="#94a3b8", label="Idle")
    axs[1].bar(labels, dyn, bottom=idle, color="#60a5fa", label="Dynamic")
    axs[1].bar(labels, wake, bottom=np.array(idle)+np.array(dyn), color="#fbbf24", label="Wake")
    axs[1].set_title("Energy Breakdown"); axs[1].set_yticks([]); axs[1].legend(frameon=False, fontsize=8)

    # 3) Heatmap (no axes ticks)
    data = np.random.rand(8,16)
    axs[2].imshow(data, aspect="auto", cmap="YlGn")
    axs[2].set_title("Efficiency Across Cells/Time")
    axs[2].set_xticks([]); axs[2].set_yticks([])

    fig_save(fig, path)

# ----------- Generate visual assets -----------
topology_png = os.path.join(ASSETS, "topology_annotated.png")
arch_png = os.path.join(ASSETS, "architecture_simple.png")
rl_png = os.path.join(ASSETS, "rl_loop.png")
energy_png = os.path.join(ASSETS, "energy_states.png")
dash_png = os.path.join(ASSETS, "dashboard_mock.png")

draw_topology_annotated(topology_png)
draw_architecture(arch_png)
draw_rl_loop(rl_png)
draw_energy_states(energy_png)
draw_dashboard_mock(dash_png)

# -------------- PPT helpers --------------
prs = Presentation()
LAY_TITLE = prs.slide_layouts[0]
LAY_TITLE_CONTENT = prs.slide_layouts[1]
LAY_BLANK = prs.slide_layouts[6]

def set_bg(slide, color=COL_BG):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color

def style_title(shape, text):
    tf = shape.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    p.font.name = FONT; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = COL_TITLE

def style_subtitle(shape, text):
    tf = shape.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    p.font.name = FONT; p.font.size = Pt(18); p.font.color.rgb = COL_TEXT

def add_bullets(slide, title, bullets):
    set_bg(slide)
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.name = FONT; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = COL_TITLE
    body = slide.placeholders[1].text_frame; body.clear()
    for i,b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b; p.level = 0
        p.font.name = FONT; p.font.size = Pt(20); p.font.color.rgb = COL_TEXT

def add_picture(slide, img, title=None, caption=None, width=Inches(8.9)):
    set_bg(slide)
    if title:
        slide.shapes.title.text = title
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.name = FONT; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = COL_TITLE
    slide.shapes.add_picture(img, Inches(0.7), Inches(1.4), width=width)
    if caption:
        tx = slide.shapes.add_textbox(Inches(0.7), Inches(5.1), Inches(8.9), Inches(0.6))
        tf = tx.text_frame; tf.text = caption
        tf.paragraphs[0].font.name = FONT; tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.color.rgb = COL_TEXT

# -------------- Build slides (10) --------------

# 1) Title
slide = prs.slides.add_slide(LAY_TITLE)
set_bg(slide)
style_title(slide.shapes.title, "Energy‑Efficient Wireless Network Simulator")
style_subtitle(slide.placeholders[1], "Visual overview: topology, architecture, RL concept, energy states, and dashboard")

# 2) Why this matters
slide = prs.slides.add_slide(LAY_TITLE_CONTENT)
add_bullets(slide, "Why this matters", [
    "Dense 5G/6G increases network energy use.",
    "Static allocation wastes power during low demand.",
    "We simulate and learn adaptive policies that save energy with QoS guardrails."
])

# 3) Concept overview
slide = prs.slides.add_slide(LAY_TITLE_CONTENT)
add_bullets(slide, "Concept overview", [
    "Realistic simulator: radio, traffic, mobility, interference.",
    "Compare baselines (PF/RR, thresholds) vs RL policy.",
    "Visual outputs: power, latency, and efficiency trends."
])

# 4) Architecture (visual)
slide = prs.slides.add_slide(LAY_TITLE_CONTENT)
add_picture(slide, arch_png, title="Architecture (modular & pluggable)",
            caption="Topology → Traffic/Mobility → Channel → MAC/PHY → Energy Model → Control → Metrics")

# 5) Topology (annotated visual)
slide = prs.slides.add_slide(LAY_TITLE_CONTENT)
add_picture(slide, topology_png, title="Topology & Assumptions",
            caption="7‑site hex (reuse‑1), mobility and interference shown. Each site has 3 sectors. RL controls sleep/power with QoS guardrails.")

# 6) Channel modeling
slide = prs.slides.add_slide(LAY_TITLE_CONTENT)
add_bullets(slide, "Channel modeling (3GPP‑style)", [
    "Path loss, shadowing, and small‑scale fading.",
    "Per‑TTI: SINR → CQI → MCS → achievable rate.",
    "Sufficient realism for energy/QoS studies without full PHY."
])

# 7) Energy model (visual)
slide = prs.slides.add_slide(LAY_TITLE_CONTENT)
add_picture(slide, energy_png, title="Energy states & transitions",
            caption="Active ↔ Micro‑sleep ↔ Deep sleep; guard timers and QoS checks before deeper sleep.")

# 8) RL loop (visual)
slide = prs.slides.add_slide(LAY_TITLE_CONTENT)
add_picture(slide, rl_png, title="Reinforcement learning loop",
            caption="Agent observes state and chooses actions (sleep/power/allocation) to reduce energy while respecting QoS.")

# 9) Dashboard (visual)
slide = prs.slides.add_slide(LAY_TITLE_CONTENT)
add_picture(slide, dash_png, title="Dashboard views",
            caption="Time series, energy breakdown, and efficiency heatmap for transparent analysis.")

# 10) Evaluation & next steps
slide = prs.slides.add_slide(LAY_TITLE_CONTENT)
add_bullets(slide, "Evaluation & next steps", [
    "Baselines: PF/RR; threshold sleep vs RL policy.",
    "Scenarios: light/moderate/peak loads; mobility present.",
    "Next: calibrate models, train policy, and compare visually.",
    "Expectation: noticeable off‑peak energy savings without QoS regressions."
])

prs.save(OUT)
print(f"Generated: {OUT}\nAssets in: {ASSETS}")